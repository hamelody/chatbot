import streamlit as st
# st.set_page_config는 반드시 Streamlit 명령어 중 가장 먼저 실행되어야 합니다.
st.set_page_config(
    page_title="유앤생명과학 업무 가이드 봇",
    layout="centered", # 또는 "wide"
    initial_sidebar_state="auto" # 또는 "expanded", "collapsed"
)

import os
import io
import fitz # PyMuPDF
import pandas as pd
import docx
from pptx import Presentation
import faiss
import openai
import numpy as np
import json
import time
from datetime import datetime
import uuid # 고유 ID 생성을 위해 추가
from openai import AzureOpenAI, APIConnectionError, APITimeoutError, RateLimitError, APIStatusError
from azure.core.exceptions import AzureError
from azure.storage.blob import BlobServiceClient
import tempfile
from werkzeug.security import check_password_hash, generate_password_hash
import traceback
import base64
import tiktoken
import re # 주석 제거 또는 다른 정규식 사용을 위해

from streamlit_cookies_manager import EncryptedCookieManager
print("Imported streamlit_cookies_manager (EncryptedCookieManager only).")


try:
    tokenizer = tiktoken.get_encoding("o200k_base") # 최신 모델용 인코더
    print("Tiktoken 'o200k_base' encoder loaded successfully.")
except Exception as e:
    st.error(f"Tiktoken encoder 'o200k_base' load failed: {e}. Token-based length limit may not work.")
    print(f"ERROR: Failed to load tiktoken 'o200k_base' encoder: {e}")
    try:
        tokenizer = tiktoken.get_encoding("cl100k_base") # 대체 인코더
        print("Tiktoken 'cl100k_base' encoder loaded successfully as a fallback.")
    except Exception as e2:
        st.error(f"Tiktoken encoder 'cl100k_base' (fallback) load failed: {e2}. Token-based length limit may not work.")
        print(f"ERROR: Failed to load tiktoken 'cl100k_base' (fallback) encoder: {e2}")
        tokenizer = None


APP_VERSION = "1.0.7 (Chat History Deletion)" 

# --- 파일 경로 및 상수 정의 ---
RULES_PATH_REPO = ".streamlit/prompt_rules.txt"
COMPANY_LOGO_PATH_REPO = "company_logo.png" # 앱 루트 디렉토리에 로고 파일 위치 가정
INDEX_BLOB_NAME = "vector_db/vector.index"
METADATA_BLOB_NAME = "vector_db/metadata.json"
USERS_BLOB_NAME = "app_data/users.json"
UPLOAD_LOG_BLOB_NAME = "app_logs/upload_log.json"
USAGE_LOG_BLOB_NAME = "app_logs/usage_log.json"
CHAT_HISTORY_BASE_PATH = "chat_histories/" # 사용자별 대화 내역 저장 기본 경로

# --- API 및 모델 설정 ---
AZURE_OPENAI_TIMEOUT = 60.0 # Azure OpenAI API 호출 타임아웃 (초)
MODEL_MAX_INPUT_TOKENS = 128000 # 사용하는 LLM의 최대 입력 토큰 수 (예: gpt-4-turbo)
MODEL_MAX_OUTPUT_TOKENS = 4096 # LLM의 최대 출력 토큰 수 (조정 가능)
BUFFER_TOKENS = 500 # 프롬프트 구성 시 여유 토큰
TARGET_INPUT_TOKENS_FOR_PROMPT = MODEL_MAX_INPUT_TOKENS - MODEL_MAX_OUTPUT_TOKENS - BUFFER_TOKENS
IMAGE_DESCRIPTION_MAX_TOKENS = 500 # 이미지 설명 생성 시 최대 토큰
EMBEDDING_BATCH_SIZE = 16 # 임베딩 배치 크기

# --- 대화 내역 관련 함수 ---
def get_current_user_login_id():
    user_info = st.session_state.get("user", {})
    return user_info.get("uid")

def get_user_chat_history_blob_name(user_login_id):
    if not user_login_id:
        return None
    return f"{CHAT_HISTORY_BASE_PATH}{user_login_id}_history.json"

def load_user_conversations_from_blob():
    user_login_id = get_current_user_login_id()
    if not user_login_id or not container_client:
        print(f"Cannot load chat history: User ID ('{user_login_id}') or container_client is missing.")
        return []
    blob_name = get_user_chat_history_blob_name(user_login_id)
    # load_data_from_blob 함수는 default_value로 {"conversations": []}와 유사한 구조를 반환하도록 수정 필요
    history_data_container = load_data_from_blob(blob_name, container_client, f"chat history for {user_login_id}", default_value={"conversations": []})
    
    # history_data_container가 None일 경우를 대비 (load_data_from_blob이 None을 반환할 수 있으므로)
    if history_data_container is None:
        loaded_conversations = []
    elif isinstance(history_data_container, dict):
        loaded_conversations = history_data_container.get("conversations", [])
    elif isinstance(history_data_container, list): # 이전 버전 호환 (list 자체를 저장한 경우)
        loaded_conversations = history_data_container
        print(f"Warning: Loaded chat history for user '{user_login_id}' as a direct list. ควรจะเป็น dict {{'conversations': [...]}}.")
    else:
        loaded_conversations = []
        print(f"Warning: Unexpected data type for chat history for user '{user_login_id}': {type(history_data_container)}")

    print(f"Loaded {len(loaded_conversations)} conversations for user '{user_login_id}'.")
    try:
        # last_updated가 없는 경우를 대비하여 get의 두 번째 인자로 기본값 제공
        loaded_conversations.sort(key=lambda x: x.get("last_updated", x.get("timestamp", "1970-01-01T00:00:00")), reverse=True)
    except Exception as e_sort:
        print(f"Error sorting conversations for user '{user_login_id}': {e_sort}")
    return loaded_conversations


def save_user_conversations_to_blob():
    user_login_id = get_current_user_login_id()
    if not user_login_id or not container_client or "all_user_conversations" not in st.session_state:
        print(f"Cannot save chat history: User ID ('{user_login_id}'), container_client, or all_user_conversations missing.")
        return False
    
    # 저장 전 최신순으로 다시 정렬 (last_updated 기준)
    try:
        st.session_state.all_user_conversations.sort(key=lambda x: x.get("last_updated", x.get("timestamp", "1970-01-01T00:00:00")), reverse=True)
    except Exception as e_sort_save:
        print(f"Error sorting conversations before saving for user '{user_login_id}': {e_sort_save}")

    blob_name = get_user_chat_history_blob_name(user_login_id)
    print(f"Saving {len(st.session_state.all_user_conversations)} conversations for user '{user_login_id}' to {blob_name}.")
    return save_data_to_blob({"conversations": st.session_state.all_user_conversations}, blob_name, container_client, f"chat history for {user_login_id}")

def generate_conversation_title(messages_list):
    if not messages_list:
        return "빈 대화"
    for msg in messages_list:
        if msg.get("role") == "user" and msg.get("content","").strip():
            # "(첨부 파일: ...)" 부분 제외
            title_candidate = msg["content"].split("\n(첨부 파일:")[0].strip()
            # 너무 길면 자르기
            return title_candidate[:30] + "..." if len(title_candidate) > 30 else title_candidate
    return "대화 시작" # 사용자 메시지가 없는 경우

def archive_current_chat_session_if_needed():
    user_login_id = get_current_user_login_id()
    # 현재 메시지가 없거나, 사용자가 없으면 아카이브할 필요 없음
    if not user_login_id or not st.session_state.get("current_chat_messages"):
        print("Archive check: No user ID or no current messages. Skipping archive.")
        return False # 변경 없음

    active_id = st.session_state.get("active_conversation_id")
    current_messages_copy = list(st.session_state.current_chat_messages) # 항상 복사본 사용
    
    archived_or_updated = False

    if active_id: # 현재 불러온 대화가 있는 경우 (업데이트 시도)
        found_and_updated = False
        for i, conv in enumerate(st.session_state.all_user_conversations):
            if conv["id"] == active_id:
                # 메시지 내용이 실제로 변경되었는지 간단히 확인 (더 정교한 비교도 가능)
                if conv["messages"] != current_messages_copy: # 메시지 목록 자체가 변경되었는지 확인
                    conv["messages"] = current_messages_copy
                    conv["last_updated"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    # 제목은 첫 메시지 기준으로 생성되었으므로, 일반적으로는 업데이트하지 않음.
                    # 필요하다면 여기서 conv["title"] = generate_conversation_title(current_messages_copy) 추가 가능
                    st.session_state.all_user_conversations[i] = conv # 리스트 내 객체 직접 수정이 반영되도록
                    print(f"Archived (updated) conversation ID: {active_id}, Title: '{conv.get('title', 'N/A')}'")
                    archived_or_updated = True
                else:
                    print(f"Conversation ID: {active_id} has no changes to messages. No update to archive needed.")
                found_and_updated = True
                break
        if not found_and_updated: # active_id가 있었지만 목록에 없는 이상한 경우 (새 대화로 처리)
             print(f"Warning: active_conversation_id '{active_id}' not found in conversation log. Treating as new chat for archiving.")
             active_id = None # 새 대화로 취급하도록 active_id 초기화
    
    # active_id가 None이거나, 위에서 None으로 바뀐 경우 (즉, 새 대화로 취급)
    if not active_id : 
        # current_chat_messages가 실제로 내용이 있어야 새 대화로 저장
        if current_messages_copy: 
            new_conv_id = str(uuid.uuid4()) # 고유 ID 생성
            title = generate_conversation_title(current_messages_copy)
            # 첫 메시지 시간 또는 현재 시간으로 대표 시간 설정
            timestamp_str = current_messages_copy[0].get("time") if current_messages_copy and current_messages_copy[0].get("time") else datetime.now().strftime("%Y-%m-%d %H:%M")

            new_conversation = {
                "id": new_conv_id,
                "title": title,
                "timestamp": timestamp_str, # 대화 시작 시점 (첫 메시지 시간 또는 생성 시간)
                "messages": current_messages_copy,
                "last_updated": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            }
            st.session_state.all_user_conversations.insert(0, new_conversation) # 최신 대화를 맨 앞에 추가
            # 새 대화가 저장되었으므로, 이제 이 대화가 "활성" 대화가 됨 (ID를 부여받았음)
            # 하지만 이 함수는 보통 컨텍스트 전환 직전에 호출되므로, 이 함수 내에서 active_id를 바꾸는 것은
            # 호출한 쪽의 로직과 꼬일 수 있음. 호출한 쪽에서 active_id를 관리하도록 둠.
            print(f"Archived (new) conversation ID: {new_conv_id}, Title: '{title}'")
            archived_or_updated = True
        else: # current_messages_copy가 비어있으면 새 대화로 저장할 내용 없음
             print("Archive check: Current messages empty and no active_id. Skipping archive of new chat.")


    if archived_or_updated: # 실제 변경/추가가 있었던 경우에만 저장
        save_user_conversations_to_blob()
    
    return archived_or_updated # 변경이 있었는지 여부 반환
# --- END 대화 내역 관련 함수 ---

def get_base64_of_bin_file(bin_file_path):
    try:
        with open(bin_file_path, 'rb') as f:
            data = f.read()
        return base64.b64encode(data).decode()
    except FileNotFoundError:
        print(f"ERROR: Logo file not found: {bin_file_path}")
        return None
    except Exception as e:
        print(f"ERROR: Processing logo file '{bin_file_path}': {e}")
        return None

def get_logo_and_version_html(app_version_str):
    logo_html_part = ""
    company_name_default = '<span class="version-text" style="font-weight:bold; font-size: 1.5em;">유앤생명과학</span>'
    
    if os.path.exists(COMPANY_LOGO_PATH_REPO):
        logo_b64 = get_base64_of_bin_file(COMPANY_LOGO_PATH_REPO)
        if logo_b64:
            logo_html_part = f'<img src="data:image/png;base64,{logo_b64}" class="logo-image" width="150" style="vertical-align: middle;">'
        else:
            logo_html_part = company_name_default
    else:
        print(f"WARNING: Company logo file not found at {COMPANY_LOGO_PATH_REPO}")
        logo_html_part = company_name_default
        
    return f"""
        {logo_html_part}
        <span class="version-text" style="vertical-align: middle; margin-left: 10px;">{app_version_str}</span>
    """

st.markdown("""
<style>
    /* 기본 CSS 스타일 */
    .stApp > header ~ div [data-testid="stHorizontalBlock"] > div:nth-child(2) div[data-testid="stButton"] > button {
        background-color: #FFFFFF; color: #333F48; border: 1px solid #BCC0C4;
        border-radius: 8px; padding: 8px 12px; font-weight: 500;
        width: auto; min-width: 100px; white-space: nowrap; display: inline-block;
        transition: background-color 0.3s ease, border-color 0.3s ease;
    }
    .stApp > header ~ div [data-testid="stHorizontalBlock"] > div:nth-child(2) div[data-testid="stButton"] > button:hover {
        background-color: #F0F2F5; border-color: #A0A4A8;
    }
    .stApp > header ~ div [data-testid="stHorizontalBlock"] > div:nth-child(2) div[data-testid="stButton"] > button:active {
        background-color: #E0E2E5;
    }
    .chat-bubble-container { display: flex; flex-direction: column; margin-bottom: 15px; }
    .user-align { align-items: flex-end; }
    .assistant-align { align-items: flex-start; }
    .bubble { padding: 10px 15px; border-radius: 18px; max-width: 75%; word-wrap: break-word; display: inline-block; color: black; box-shadow: 0px 2px 4px rgba(0, 0, 0, 0.08); position: relative; }
    .user-bubble { background-color: #90EE90; color: black; border-bottom-right-radius: 5px; }
    .assistant-bubble { background-color: #E9E9EB; border-bottom-left-radius: 5px; }
    .timestamp { font-size: 0.7rem; color: #8E8E93; padding: 2px 5px 0px 5px; }
    
    /* 메인 앱 제목 및 부제목 (로그인 후) */
    .main-app-title-container { text-align: center; margin-bottom: 24px; }
    .main-app-title { font-size: 2.1rem; font-weight: bold; display: block; }
    .main-app-subtitle { font-size: 0.9rem; color: gray; display: block; margin-top: 4px;}
    
    /* 로고 및 버전 */
    .logo-container { display: flex; align-items: center; }
    .logo-image { margin-right: 10px; }
    .version-text { font-size: 0.9rem; color: gray; }

    /* 로그인 화면 전용 제목 스타일 */
    .login-page-header-container { text-align: center; margin-top: 10px; margin-bottom: 10px;}
    .login-page-main-title { font-size: 1.8rem; font-weight: bold; display: block; color: #333F48; } 
    .login-page-sub-title { font-size: 0.85rem; color: gray; display: block; margin-top: 2px; margin-bottom: 20px;}
    .login-form-title { 
        font-size: 1.6rem; 
        font-weight: bold;
        text-align: center;
        margin-top: 10px; 
        margin-bottom: 25px; 
    }

    /* 모바일 화면 대응 */
    @media (max-width: 768px) {
        .main-app-title { font-size: 1.8rem; }
        .main-app-subtitle { font-size: 0.8rem; }
        .login-page-main-title { font-size: 1.5rem; }
        .login-page-sub-title { font-size: 0.8rem; }
        .login-form-title { font-size: 1.3rem; margin-bottom: 20px; }
        .stButton>button { font-size: 0.9rem !important; } /* 사이드바 버튼 등 모바일 크기 조정 */
    }
    /* 사이드바 버튼 스타일 미세 조정 */
    .stButton>button { 
        text-align: left !important; 
        display: block !important; 
        white-space: nowrap;
        overflow: hidden;
        text-overflow: ellipsis;
    }
</style>
""", unsafe_allow_html=True)

@st.cache_resource
def get_azure_openai_client_cached():
    print("Attempting to initialize Azure OpenAI client...")
    try:
        api_version_to_use = st.secrets.get("AZURE_OPENAI_VERSION", "2024-02-15-preview")
        print(f"DEBUG: Initializing AzureOpenAI client with API version: {api_version_to_use}")
        client = AzureOpenAI(
            api_key=st.secrets["AZURE_OPENAI_KEY"],
            azure_endpoint=st.secrets["AZURE_OPENAI_ENDPOINT"],
            api_version=api_version_to_use,
            timeout=AZURE_OPENAI_TIMEOUT
        )
        print("Azure OpenAI client initialized successfully.")
        return client
    except KeyError as e:
        st.error(f"Azure OpenAI config error: Missing key '{e.args[0]}' in secrets.")
        print(f"ERROR: Missing Azure OpenAI secret: {e.args[0]}")
        return None
    except Exception as e:
        st.error(f"Error initializing Azure OpenAI client: {e}.")
        print(f"ERROR: Azure OpenAI client initialization failed: {e}\n{traceback.format_exc()}")
        return None

@st.cache_resource
def get_azure_blob_clients_cached():
    print("Attempting to initialize Azure Blob Service client...")
    try:
        conn_str = st.secrets["AZURE_BLOB_CONN"]
        blob_service_client = BlobServiceClient.from_connection_string(conn_str, connection_timeout=60, read_timeout=120)
        container_name = st.secrets["BLOB_CONTAINER"]
        container_client = blob_service_client.get_container_client(container_name)
        print(f"Azure Blob Service client and container client for '{container_name}' initialized successfully.")
        return blob_service_client, container_client
    except KeyError as e:
        st.error(f"Azure Blob Storage config error: Missing key '{e.args[0]}' in secrets.")
        print(f"ERROR: Missing Azure Blob Storage secret: {e.args[0]}")
        return None, None
    except Exception as e:
        st.error(f"Error initializing Azure Blob client: {e}.")
        print(f"ERROR: Azure Blob client initialization failed: {e}\n{traceback.format_exc()}")
        return None, None

openai_client = get_azure_openai_client_cached()
blob_service, container_client = get_azure_blob_clients_cached()

EMBEDDING_MODEL = None
if openai_client:
    try:
        EMBEDDING_MODEL = st.secrets["AZURE_OPENAI_EMBEDDING_DEPLOYMENT"]
        print(f"Embedding model set to: {EMBEDDING_MODEL}")
    except KeyError:
        st.error("Missing 'AZURE_OPENAI_EMBEDDING_DEPLOYMENT' in secrets.")
        print("ERROR: Missing AZURE_OPENAI_EMBEDDING_DEPLOYMENT secret.")
        openai_client = None # 클라이언트 사용 불가 처리
    except Exception as e:
        st.error(f"Error loading embedding model config: {e}")
        print(f"ERROR: Loading AZURE_OPENAI_EMBEDDING_DEPLOYMENT secret: {e}")
        openai_client = None # 클라이언트 사용 불가 처리


def load_data_from_blob(blob_name, _container_client, data_description="data", default_value=None):
    if not _container_client:
        print(f"ERROR: Blob Container client is None for load_data_from_blob ('{data_description}'). Returning default.")
        return default_value if default_value is not None else ({} if not isinstance(default_value, list) else [])
    
    print(f"Attempting to load '{data_description}' from Blob: '{blob_name}'")
    try:
        blob_client_instance = _container_client.get_blob_client(blob_name)
        if blob_client_instance.exists():
            with tempfile.TemporaryDirectory() as tmpdir:
                local_temp_path = os.path.join(tmpdir, os.path.basename(blob_name))
                with open(local_temp_path, "wb") as download_file:
                    download_stream = blob_client_instance.download_blob(timeout=60)
                    download_file.write(download_stream.readall())
                
                if os.path.getsize(local_temp_path) > 0:
                    with open(local_temp_path, "r", encoding="utf-8") as f:
                        loaded_data = json.load(f)
                    print(f"Successfully loaded '{data_description}' from Blob: '{blob_name}'")
                    return loaded_data
                else: # 파일은 존재하나 비어있는 경우
                    print(f"WARNING: '{data_description}' file '{blob_name}' exists in Blob but is empty. Returning default.")
                    return default_value if default_value is not None else ({} if not isinstance(default_value, list) else [])
        else: # 파일이 존재하지 않는 경우
            print(f"WARNING: '{data_description}' file '{blob_name}' not found in Blob Storage. Returning default.")
            return default_value if default_value is not None else ({} if not isinstance(default_value, list) else [])
    except json.JSONDecodeError: # JSON 파싱 오류
        print(f"ERROR: Failed to decode JSON for '{data_description}' from Blob '{blob_name}'. Returning default.")
        st.warning(f"File '{data_description}' ({blob_name}) is corrupted or not valid JSON. Using default.")
        return default_value if default_value is not None else ({} if not isinstance(default_value, list) else [])
    except AzureError as ae: # Azure 관련 오류
        print(f"AZURE ERROR loading '{data_description}' from Blob '{blob_name}': {ae}\n{traceback.format_exc()}")
        st.warning(f"Azure service error loading '{data_description}': {ae}. Using default.")
        return default_value if default_value is not None else ({} if not isinstance(default_value, list) else [])
    except Exception as e: # 기타 모든 오류
        print(f"GENERAL ERROR loading '{data_description}' from Blob '{blob_name}': {e}\n{traceback.format_exc()}")
        st.warning(f"Unknown error loading '{data_description}': {e}. Using default.")
        return default_value if default_value is not None else ({} if not isinstance(default_value, list) else [])

def save_data_to_blob(data_to_save, blob_name, _container_client, data_description="data"):
    if not _container_client:
        # st.error(f"Cannot save '{data_description}': Azure Blob client not ready.") # UI 오류 최소화
        print(f"ERROR: Blob Container client is None, cannot save '{data_description}' to '{blob_name}'.")
        return False
    try:
        # JSON 직렬화 가능한 타입인지 확인 (dict 또는 list)
        if not isinstance(data_to_save, (dict, list)):
            # st.error(f"Save failed for '{data_description}': Data is not JSON serializable (type: {type(data_to_save)}).")
            print(f"ERROR: Data for '{blob_name}' is not JSON serializable (type: {type(data_to_save)}).")
            return False
            
        with tempfile.TemporaryDirectory() as tmpdir:
            local_temp_path = os.path.join(tmpdir, os.path.basename(blob_name))
            with open(local_temp_path, "w", encoding="utf-8") as f:
                json.dump(data_to_save, f, ensure_ascii=False, indent=2)
            
            blob_client_instance = _container_client.get_blob_client(blob_name)
            with open(local_temp_path, "rb") as data_stream:
                blob_client_instance.upload_blob(data_stream, overwrite=True, timeout=60)
            print(f"Successfully saved '{data_description}' to Blob: '{blob_name}'")
        return True
    except AzureError as ae:
        # st.error(f"Azure service error saving '{data_description}' to Blob: {ae}")
        print(f"AZURE ERROR saving '{data_description}' to Blob '{blob_name}': {ae}\n{traceback.format_exc()}")
        return False
    except Exception as e:
        # st.error(f"Unknown error saving '{data_description}' to Blob: {e}")
        print(f"GENERAL ERROR saving '{data_description}' to Blob '{blob_name}': {e}\n{traceback.format_exc()}")
        return False

def save_binary_data_to_blob(local_file_path, blob_name, _container_client, data_description="binary data"):
    if not _container_client:
        # st.error(f"Cannot save binary '{data_description}': Azure Blob client not ready.")
        print(f"ERROR: Blob Container client is None, cannot save binary '{blob_name}'.")
        return False
    if not os.path.exists(local_file_path):
        # st.error(f"Local file for binary '{data_description}' not found: '{local_file_path}'")
        print(f"ERROR: Local file for binary data not found: '{local_file_path}'")
        return False
    try:
        blob_client_instance = _container_client.get_blob_client(blob_name)
        with open(local_file_path, "rb") as data_stream:
            blob_client_instance.upload_blob(data_stream, overwrite=True, timeout=120) # 바이너리 파일은 타임아웃 길게
        print(f"Successfully saved binary '{data_description}' to Blob: '{blob_name}'")
        return True
    except AzureError as ae:
        # st.error(f"Azure service error saving binary '{data_description}' to Blob: {ae}")
        print(f"AZURE ERROR saving binary '{data_description}' to Blob '{blob_name}': {ae}\n{traceback.format_exc()}")
        return False
    except Exception as e:
        # st.error(f"Unknown error saving binary '{data_description}' to Blob: {e}")
        print(f"GENERAL ERROR saving binary '{data_description}' to Blob '{blob_name}': {e}\n{traceback.format_exc()}")
        return False

USERS = {}
if container_client:
    USERS = load_data_from_blob(USERS_BLOB_NAME, container_client, "user info", default_value={})
    if not isinstance(USERS, dict) : # 로드된 데이터가 dict가 아니면 초기화
        print(f"ERROR: USERS loaded from blob is not a dict ({type(USERS)}). Re-initializing.")
        USERS = {}
    if "admin" not in USERS: # admin 계정이 없으면 기본값으로 생성
        print(f"'{USERS_BLOB_NAME}' from Blob is empty or admin is missing. Creating default admin.")
        admin_password = st.secrets.get("ADMIN_PASSWORD", "diteam_fallback_secret") # ADMIN_PASSWORD secrets에서 가져오기
        USERS["admin"] = {
            "name": "관리자", "department": "품질보증팀", "uid": "admin", # uid 필드 추가
            "password_hash": generate_password_hash(admin_password),
            "approved": True, "role": "admin"
        }
        if not save_data_to_blob(USERS, USERS_BLOB_NAME, container_client, "initial user info with default admin"):
             st.warning("Failed to save default admin info to Blob. Will retry on next user data save.") # UI 경고
else: # Blob 클라이언트 연결 실패 시
    st.error("Azure Blob Storage connection failed. Cannot initialize user info. App may not function correctly.")
    print("CRITICAL: Cannot initialize USERS due to Blob client failure.")
    USERS = {"admin": {"name": "관리자(연결실패)", "department": "시스템", "uid":"admin", "password_hash": generate_password_hash("fallback"), "approved": True, "role": "admin"}}


cookies = None
cookie_manager_ready = False # 전역 변수처럼 사용될 쿠키 매니저 준비 상태
print(f"Attempting to load COOKIE_SECRET from st.secrets...")
try:
    cookie_secret_key = st.secrets.get("COOKIE_SECRET")
    if not cookie_secret_key:
        # st.error("'COOKIE_SECRET' is not set or empty in st.secrets.") # 로그인 전 UI 오류 최소화
        print("ERROR: COOKIE_SECRET is not set or empty in st.secrets.")
    else:
        cookies = EncryptedCookieManager(
            prefix="gmp_chatbot_v1.0.7_cookie/", # 쿠키 prefix 변경 (버전업 및 기능명시)
            password=cookie_secret_key
        )
        print("CookieManager object created. Readiness will be checked before use.")
except Exception as e:
    # st.error(f"Unknown error creating cookie manager object: {e}") # 로그인 전 UI 오류 최소화
    print(f"CRITICAL: CookieManager object creation error: {e}\n{traceback.format_exc()}")
    cookies = None # 쿠키 객체 생성 실패 시 None으로 설정

SESSION_TIMEOUT = 1800 # 세션 타임아웃 기본값 (30분)
try:
    session_timeout_secret = st.secrets.get("SESSION_TIMEOUT")
    if session_timeout_secret: SESSION_TIMEOUT = int(session_timeout_secret)
    print(f"Session timeout set to: {SESSION_TIMEOUT} seconds.")
except (ValueError, TypeError):
    print(f"WARNING: SESSION_TIMEOUT in secrets ('{session_timeout_secret}') is not a valid integer. Using default {SESSION_TIMEOUT}s.")
except Exception as e:
     print(f"WARNING: Error reading SESSION_TIMEOUT from secrets: {e}. Using default {SESSION_TIMEOUT}s.")

# --- Session State 초기화 ---
session_keys_defaults = {
    "authenticated": False, "user": {},
    "current_chat_messages": [], "all_user_conversations": [],
    "active_conversation_id": None, "show_uploader": False,
    "pending_delete_conv_id": None # 대화 삭제 확인용 ID 저장
}
for key, default_value in session_keys_defaults.items():
    if key not in st.session_state:
        st.session_state[key] = default_value
        print(f"Initializing st.session_state['{key}'] to default.")

# --- 쿠키를 사용한 세션 복원 시도 (앱 실행 초기) ---
if not st.session_state.get("authenticated", False) and cookies is not None:
    print("Attempting initial session restore from cookies as user is not authenticated in session_state.")
    try:
        if cookies.ready(): # 이 시점에 ready()가 True여야 함
            cookie_manager_ready = True # True로 설정!
            print("CookieManager.ready() is True for initial session restore.")
            auth_cookie_val = cookies.get("authenticated")
            print(f"Cookie 'authenticated' value for initial restore: {auth_cookie_val}")

            if auth_cookie_val == "true":
                login_time_str = cookies.get("login_time", "0")
                try: login_time = float(login_time_str if login_time_str and login_time_str.replace('.', '', 1).isdigit() else "0")
                except ValueError: login_time = 0.0
                
                if (time.time() - login_time) < SESSION_TIMEOUT:
                    user_json_cookie = cookies.get("user", "{}")
                    try:
                        user_data_from_cookie = json.loads(user_json_cookie if user_json_cookie else "{}")
                        # uid가 있는지 확인 (중요)
                        if user_data_from_cookie and isinstance(user_data_from_cookie, dict) and "uid" in user_data_from_cookie:
                            st.session_state["user"] = user_data_from_cookie
                            st.session_state["authenticated"] = True
                            # 로그인 성공 시 대화 내역 로드
                            st.session_state.all_user_conversations = load_user_conversations_from_blob() # user_id는 내부적으로 get_current_user_login_id() 사용
                            st.session_state.current_chat_messages = [] # 새 대화로 시작
                            st.session_state.active_conversation_id = None
                            print(f"User '{user_data_from_cookie.get('name')}' session restored from cookie. Chat history loaded.")
                            # 여기서 st.rerun()을 호출하면 쿠키 관련 컴포넌트가 아직 완전히 마운트되지 않아 오류 발생 가능성 있음
                        else:
                            print("User data in cookie is empty, invalid, or missing uid. Clearing auth state from cookie.")
                            if cookies.ready(): cookies["authenticated"] = "false"; cookies["user"] = ""; cookies["login_time"] = ""; cookies.save()
                            st.session_state["authenticated"] = False
                    except json.JSONDecodeError:
                        print("ERROR: Failed to decode user JSON from cookie. Clearing auth state from cookie.")
                        if cookies.ready(): cookies["authenticated"] = "false"; cookies["user"] = ""; cookies["login_time"] = ""; cookies.save()
                        st.session_state["authenticated"] = False
                else: # 세션 타임아웃
                    print("Session timeout detected from cookie. Clearing auth state and cookies.")
                    if cookies.ready(): cookies["authenticated"] = "false"; cookies["user"] = ""; cookies["login_time"] = ""; cookies.save()
                    st.session_state["authenticated"] = False
            # else: auth_cookie_val이 "true"가 아니면 st.session_state["authenticated"]는 False로 유지됨
        else: # cookies.ready() is False
            print("CookieManager.ready() is False for initial session restore. Cannot load cookies at this exact moment.")
            # cookie_manager_ready는 False로 유지됨
    except Exception as e_cookie_op_initial:
        print(f"Exception during initial cookie operations (session restore): {e_cookie_op_initial}\n{traceback.format_exc()}")
        st.session_state["authenticated"] = False # 안전하게 False로
        # cookie_manager_ready는 False로 유지될 수 있음

# 로그인 UI 표시 전 쿠키 매니저 준비 상태 최종 확인 (위에서 ready가 아니었을 수 있으므로)
if cookies is not None and not cookie_manager_ready: # cookie_manager_ready가 여전히 False이면
    print("Checking CookieManager readiness again before login UI (if it was not ready initially)...")
    try:
        if cookies.ready():
            cookie_manager_ready = True # True로 업데이트
            print("CookieManager became ready just before login UI (second check).")
            # 만약 여기서 ready가 되었고, 아직 인증 안된 상태라면, 위 세션 복원 로직을 한번 더 시도해볼 수 있음.
            # 하지만 복잡성을 줄이기 위해, 현재는 이 플래그만 업데이트하고, 다음번 rerun 시 위 로직이 다시 시도되도록 함.
            # 또는, 여기서 명시적으로 세션 복원을 다시 시도하고 성공 시 st.rerun().
            # (1.0.5 버전의 "두 번째 기회" 로직이 이와 유사했음. 여기서는 우선 플래그 업데이트에 집중)
        else:
            print("CookieManager still not ready before login UI (second check).")
    except Exception as e_ready_login_ui:
        print(f"WARNING: cookies.ready() call just before login UI failed: {e_ready_login_ui}")


if not st.session_state.get("authenticated", False):
    # 로그인 화면 UI
    st.markdown("""
    <div class="login-page-header-container" style="margin-top: 80px;"> 
      <span class="login-page-main-title">유앤생명과학 GMP/SOP 업무 가이드 봇</span>
      <span class="login-page-sub-title">Made by DI.PART</span>
    </div>
    """, unsafe_allow_html=True)
    st.markdown('<p class="login-form-title">🔐 로그인 또는 회원가입</p>', unsafe_allow_html=True)

    if cookies is None or not cookie_manager_ready: # 최종적으로 쿠키 매니저 상태 확인 후 경고
        st.warning("쿠키 시스템이 아직 준비되지 않았습니다. 로그인이 유지되지 않을 수 있습니다. 잠시 후 새로고침 해보세요.")

    with st.form("auth_form_v1.0.7_delete", clear_on_submit=False):
        mode = st.radio("선택", ["로그인", "회원가입"], key="auth_mode_v1.0.7_delete")
        # form 내부에서는 입력값을 변수에 할당해도 submit 전까지는 외부에서 사용 불가
        uid_input_form = st.text_input("ID", key="auth_uid_v1.0.7_delete")
        pwd_form = st.text_input("비밀번호", type="password", key="auth_pwd_v1.0.7_delete")
        name_form, dept_form = "", ""
        if mode == "회원가입":
            name_form = st.text_input("이름", key="auth_name_v1.0.7_delete")
            dept_form = st.text_input("부서", key="auth_dept_v1.0.7_delete")
        submit_button_form = st.form_submit_button("확인")

    if submit_button_form: # 폼 제출 시에만 아래 로직 실행
        if not uid_input_form or not pwd_form: st.error("ID와 비밀번호를 모두 입력해주세요.")
        elif mode == "회원가입" and (not name_form or not dept_form): st.error("회원가입 시 이름과 부서를 모두 입력해주세요.")
        else:
            if mode == "로그인":
                user_data_from_db = USERS.get(uid_input_form) # DB(USERS 딕셔너리)에서 사용자 정보 가져오기
                if not user_data_from_db: st.error("존재하지 않는 ID입니다.")
                elif not user_data_from_db.get("approved", False): st.warning("관리자 승인 대기 중인 계정입니다.")
                elif check_password_hash(user_data_from_db["password_hash"], pwd_form):
                    
                    # 세션에 저장할 사용자 정보 구성 (uid 포함)
                    session_user_data_on_login = user_data_from_db.copy()
                    session_user_data_on_login["uid"] = uid_input_form # USERS 딕셔너리의 키(로그인 ID)를 uid로 저장

                    st.session_state["authenticated"] = True
                    st.session_state["user"] = session_user_data_on_login
                    
                    # 로그인 성공 시 대화 내역 로드 및 새 대화 준비
                    st.session_state.all_user_conversations = load_user_conversations_from_blob() # user_id는 내부적으로 get_current_user_login_id() 사용
                    st.session_state.current_chat_messages = [] # 새 대화로 시작
                    st.session_state.active_conversation_id = None
                    st.session_state.pending_delete_conv_id = None # 혹시 남아있을 수 있는 플래그 초기화
                    print(f"Login successful for user '{uid_input_form}'. Chat history loaded. Starting new chat session.")

                    if cookies is not None and cookie_manager_ready: # 쿠키 매니저 준비되었을 때만 쿠키 저장
                        try:
                            cookies["authenticated"] = "true"
                            cookies["user"] = json.dumps(session_user_data_on_login) # uid 포함된 정보 저장
                            cookies["login_time"] = str(time.time())
                            cookies.save()
                            print(f"Cookies saved for user '{uid_input_form}'.")
                        except Exception as e_cookie_save_login:
                            st.warning(f"로그인 쿠키 저장 중 문제 발생: {e_cookie_save_login}")
                            print(f"ERROR: Failed to save login cookies: {e_cookie_save_login}")
                    elif cookies is None:
                         st.warning("쿠키 시스템이 초기화되지 않아 로그인 상태를 저장할 수 없습니다. (로그인 시점)")
                    elif not cookie_manager_ready:
                         st.warning("쿠키 시스템이 준비되지 않아 로그인 상태를 저장할 수 없습니다. (로그인 시점, not ready)")
                    
                    st.success(f"{session_user_data_on_login.get('name', uid_input_form)}님, 환영합니다!"); st.rerun()
                else: st.error("비밀번호가 일치하지 않습니다.")
            elif mode == "회원가입":
                if uid_input_form in USERS: st.error("이미 존재하는 ID입니다.")
                else:
                    USERS[uid_input_form] = {"name": name_form, "department": dept_form, "uid": uid_input_form, 
                                  "password_hash": generate_password_hash(pwd_form),
                                  "approved": False, "role": "user"}
                    if not save_data_to_blob(USERS, USERS_BLOB_NAME, container_client, "user info (signup)"):
                        st.error("회원 정보 저장에 실패했습니다. 관리자에게 문의하세요.")
                        USERS.pop(uid_input_form, None) # 저장 실패 시 롤백
                    else:
                        st.success("회원가입 요청이 완료되었습니다. 관리자 승인 후 로그인 가능합니다.")
    st.stop() # 인증되지 않은 사용자는 여기서 실행 중지

# --- 이하 코드는 인증된 사용자에게만 보임 ---
current_user_info = st.session_state.get("user", {}) # uid 포함

# --- 사이드바: 사용자 정보, 새 대화 버튼 및 대화 내역 ---
with st.sidebar:
    st.markdown(f"**{current_user_info.get('name', '사용자')}** (`{current_user_info.get('uid', 'ID없음')}`)")
    st.markdown(f"*{current_user_info.get('department', '부서정보없음')}*")
    st.markdown("---")

    if st.button("➕ 새 대화 시작", use_container_width=True, key="new_chat_button_sidebar_v7"):
        # 현재 대화가 있다면 아카이브 (ID가 없어도 새 대화로 아카이브됨)
        archive_current_chat_session_if_needed() 
        
        st.session_state.current_chat_messages = [] # 현재 채팅 메시지 비우기
        st.session_state.active_conversation_id = None # 활성 대화 ID 없음 (새 대화 상태)
        st.session_state.pending_delete_conv_id = None # 삭제 보류 ID 초기화
        print("New chat started by user via sidebar button.")
        st.rerun()

    st.markdown("##### 이전 대화")
    
    # 삭제 확인 UI (pending_delete_conv_id가 있을 때만 표시)
    if st.session_state.get("pending_delete_conv_id"):
        conv_id_to_delete = st.session_state.pending_delete_conv_id
        # 삭제할 대화의 제목 찾기 (없으면 ID로 대체)
        conv_title_to_delete = conv_id_to_delete 
        for c_del in st.session_state.all_user_conversations:
            if c_del['id'] == conv_id_to_delete:
                conv_title_to_delete = c_del.get('title', conv_id_to_delete)
                break
        
        st.sidebar.warning(f"'{conv_title_to_delete}' 대화를 정말 삭제하시겠습니까?")
        del_confirm_cols = st.sidebar.columns(2)
        if del_confirm_cols[0].button("✅ 예, 삭제", key=f"confirm_del_yes_{conv_id_to_delete}", use_container_width=True):
            # all_user_conversations 리스트에서 해당 ID의 대화 제거
            st.session_state.all_user_conversations = [
                c for c in st.session_state.all_user_conversations if c['id'] != conv_id_to_delete
            ]
            save_user_conversations_to_blob() # 변경된 전체 목록 저장
            
            # 만약 현재 활성 대화가 삭제된 대화였다면, 현재 채팅창 비우기
            if st.session_state.active_conversation_id == conv_id_to_delete:
                st.session_state.current_chat_messages = []
                st.session_state.active_conversation_id = None
            
            st.session_state.pending_delete_conv_id = None # 삭제 보류 플래그 해제
            st.toast(f"'{conv_title_to_delete}' 대화가 삭제되었습니다.", icon="🗑️")
            st.rerun()

        if del_confirm_cols[1].button("❌ 아니요", key=f"confirm_del_no_{conv_id_to_delete}", use_container_width=True):
            st.session_state.pending_delete_conv_id = None # 삭제 보류 플래그 해제
            st.rerun()

    # 대화 목록 표시
    if not st.session_state.all_user_conversations:
        st.sidebar.caption("이전 대화 내역이 없습니다.")
    
    # all_user_conversations는 load/save 시 이미 last_updated 기준 내림차순 정렬됨
    # 화면에는 최근 20개 또는 설정한 개수만큼 표시
    for conv_idx, conv_data in enumerate(st.session_state.all_user_conversations[:20]): 
        # 각 대화 아이템을 가로로 배치 (제목/시간 버튼, 삭제 아이콘 버튼)
        item_cols = st.sidebar.columns([0.85, 0.15]) # 버튼과 아이콘 비율
        
        title_display = conv_data.get('title', f"대화_{conv_data['id'][:8]}")
        # timestamp는 첫 메시지 시간, last_updated는 마지막 수정 시간. 둘 다 없으면 빈 문자열
        timestamp_display = conv_data.get('last_updated', conv_data.get('timestamp','')) 
        
        button_label = f"{title_display} ({timestamp_display})"
        is_active_conversation = (st.session_state.active_conversation_id == conv_data["id"])

        # 대화 불러오기 버튼
        if item_cols[0].button(
            button_label, 
            key=f"load_conv_btn_{conv_data['id']}", 
            use_container_width=True, 
            type="primary" if is_active_conversation else "secondary",
            help=f"{title_display} 대화 보기"
        ):
            if not is_active_conversation: # 현재 활성 대화가 아닐 때만 새로 로드
                # 현재 진행중이던 대화(current_chat_messages)가 새 내용이면 저장
                archive_current_chat_session_if_needed() 
                
                st.session_state.current_chat_messages = list(conv_data["messages"]) # 대화 내용 불러오기 (복사본)
                st.session_state.active_conversation_id = conv_data["id"]
                st.session_state.pending_delete_conv_id = None # 다른 대화 선택 시 삭제 보류 해제
                print(f"Loaded conversation ID: {conv_data['id']}, Title: '{title_display}'")
                st.rerun()
        
        # 삭제 아이콘 버튼
        if item_cols[1].button("🗑️", key=f"delete_icon_btn_{conv_data['id']}", help="이 대화 삭제"):
            st.session_state.pending_delete_conv_id = conv_data['id'] # 삭제 보류 상태로 설정
            st.rerun() # 삭제 확인 UI를 표시하기 위해 rerun

    if len(st.session_state.all_user_conversations) > 20:
        st.sidebar.caption("더 많은 내역은 전체 보기 기능(추후 구현)을 이용해주세요.")


# --- 메인 화면 상단 로고 및 로그아웃 버튼 ---
top_cols_main = st.columns([0.7, 0.3])
with top_cols_main[0]:
    main_logo_html = get_logo_and_version_html(APP_VERSION)
    st.markdown(f"""<div class="logo-container">{main_logo_html}</div>""", unsafe_allow_html=True)

with top_cols_main[1]:
    st.markdown('<div style="text-align: right;">', unsafe_allow_html=True)
    if st.button("로그아웃", key="logout_button_v1.0.7_delete"):
        archive_current_chat_session_if_needed() # 로그아웃 전 현재 대화 저장
        
        # 모든 세션 상태 초기화 (대화 내역 관련 포함)
        for key_to_reset in session_keys_defaults.keys():
            if key_to_reset == "authenticated": # authenticated는 False로
                st.session_state[key_to_reset] = False
            elif key_to_reset == "user": # user는 빈 dict로
                st.session_state[key_to_reset] = {}
            elif isinstance(session_keys_defaults[key_to_reset], list): # 리스트 타입은 빈 리스트로
                st.session_state[key_to_reset] = []
            else: # 나머지는 None 또는 기본값 (pending_delete_conv_id 등)
                st.session_state[key_to_reset] = session_keys_defaults[key_to_reset]
        
        print("Logout successful. All relevant session states cleared.")
        
        if cookies is not None and cookie_manager_ready: # 쿠키 매니저 준비된 경우에만
            try:
                cookies["authenticated"] = "false"
                cookies["user"] = ""
                cookies["login_time"] = ""
                cookies.save()
                print("Cookies cleared on logout.")
            except Exception as e_logout_cookie:
                 print(f"ERROR: Failed to clear cookies on logout: {e_logout_cookie}")
        st.rerun()
    st.markdown('</div>', unsafe_allow_html=True)

# --- 메인 화면 앱 제목 ---
st.markdown("""
<div class="main-app-title-container">
  <span class="main-app-title">유앤생명과학 GMP/SOP 업무 가이드 봇</span>
  <span class="main-app-subtitle">Made by DI.PART</span>
</div>
""", unsafe_allow_html=True)


# --- @st.cache_resource 및 @st.cache_data 함수들 ---
@st.cache_resource
def load_vector_db_from_blob_cached(_container_client):
    if not _container_client:
        print("ERROR: Blob Container client is None for load_vector_db_from_blob_cached.")
        return faiss.IndexFlatL2(1536), []
    current_embedding_dimension = 1536
    idx, meta = faiss.IndexFlatL2(current_embedding_dimension), []
    print(f"Attempting to load vector DB from Blob: '{INDEX_BLOB_NAME}', '{METADATA_BLOB_NAME}' with dimension {current_embedding_dimension}")
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            local_index_path = os.path.join(tmpdir, os.path.basename(INDEX_BLOB_NAME))
            local_metadata_path = os.path.join(tmpdir, os.path.basename(METADATA_BLOB_NAME))

            index_blob_client = _container_client.get_blob_client(INDEX_BLOB_NAME)
            if index_blob_client.exists():
                print(f"Downloading '{INDEX_BLOB_NAME}'...")
                with open(local_index_path, "wb") as download_file:
                    download_stream = index_blob_client.download_blob(timeout=60)
                    download_file.write(download_stream.readall())
                if os.path.getsize(local_index_path) > 0:
                    try:
                        idx = faiss.read_index(local_index_path)
                        if idx.d != current_embedding_dimension:
                            print(f"WARNING: Loaded FAISS index dimension ({idx.d}) does not match expected dimension ({current_embedding_dimension}). Re-initializing.")
                            idx = faiss.IndexFlatL2(current_embedding_dimension); meta = []
                        else:
                            print(f"'{INDEX_BLOB_NAME}' loaded successfully from Blob Storage. Dimension: {idx.d}")
                    except Exception as e_faiss_read:
                        print(f"ERROR reading FAISS index: {e_faiss_read}. Re-initializing index.")
                        idx = faiss.IndexFlatL2(current_embedding_dimension); meta = []
                else:
                    print(f"WARNING: '{INDEX_BLOB_NAME}' is empty in Blob. Using new index."); idx = faiss.IndexFlatL2(current_embedding_dimension); meta = []
            else:
                print(f"WARNING: '{INDEX_BLOB_NAME}' not found in Blob Storage. New index will be used/created."); idx = faiss.IndexFlatL2(current_embedding_dimension); meta = []

            if idx is not None: # idx가 성공적으로 초기화/로드 된 경우
                metadata_blob_client = _container_client.get_blob_client(METADATA_BLOB_NAME)
                # 메타데이터는 인덱스 파일이 실제로 존재하고 내용이 있거나, DB에 아이템이 있을 때만 로드 시도
                if metadata_blob_client.exists() and (idx.ntotal > 0 or (index_blob_client.exists() and os.path.exists(local_index_path) and os.path.getsize(local_index_path) > 0) ):
                    print(f"Downloading '{METADATA_BLOB_NAME}'...")
                    with open(local_metadata_path, "wb") as download_file_meta:
                        download_stream_meta = metadata_blob_client.download_blob(timeout=60)
                        download_file_meta.write(download_stream_meta.readall())
                    if os.path.getsize(local_metadata_path) > 0 :
                        with open(local_metadata_path, "r", encoding="utf-8") as f_meta: meta = json.load(f_meta)
                    else: meta = []; print(f"WARNING: '{METADATA_BLOB_NAME}' is empty in Blob.")
                # 인덱스가 새롭고 비어있으며, 인덱스 파일도 없는 경우 (완전 초기 상태)
                elif idx.ntotal == 0 and not (index_blob_client.exists() and os.path.exists(local_index_path) and os.path.getsize(local_index_path) > 0):
                     print(f"INFO: Index is new and empty, and no existing index file in blob. Starting with empty metadata."); meta = []
                else: # 메타데이터 파일이 없거나, 인덱스는 있지만 해당 인덱스 파일이 없는 등의 그 외 상황
                    print(f"INFO: Metadata file '{METADATA_BLOB_NAME}' not found, or index is empty/inconsistent with file. Starting with empty metadata."); meta = []

            # 데이터 일관성 최종 체크
            if idx is not None and idx.ntotal == 0 and len(meta) > 0: # 인덱스는 비었는데 메타데이터만 있는 경우
                print(f"INFO: FAISS index is empty (ntotal=0) but metadata is not. Clearing metadata for consistency."); meta = []
            elif idx is not None and idx.ntotal > 0 and not meta and (index_blob_client.exists() and os.path.exists(local_index_path) and os.path.getsize(local_index_path) > 0) : # 인덱스는 있는데 메타데이터가 없는 경우 (파일은 존재)
                print(f"CRITICAL WARNING: FAISS index has data (ntotal={idx.ntotal}) but metadata is empty, despite index file existing. This may lead to errors.")
    except AzureError as ae:
        st.error(f"Azure service error loading vector DB from Blob: {ae}"); print(f"AZURE ERROR loading vector DB: {ae}\n{traceback.format_exc()}"); idx = faiss.IndexFlatL2(current_embedding_dimension); meta = []
    except Exception as e:
        st.error(f"Unknown error loading vector DB from Blob: {e}"); print(f"GENERAL ERROR loading vector DB: {e}\n{traceback.format_exc()}"); idx = faiss.IndexFlatL2(current_embedding_dimension); meta = []
    return idx, meta

index, metadata = (faiss.IndexFlatL2(1536), []) # 기본값으로 초기화
if container_client: # Blob 클라이언트가 성공적으로 초기화된 경우에만 로드 시도
    index, metadata = load_vector_db_from_blob_cached(container_client)
    print(f"DEBUG: FAISS index loaded after cache. ntotal: {index.ntotal if index else 'Index is None'}, dimension: {index.d if index else 'N/A'}")
    print(f"DEBUG: Metadata loaded after cache. Length: {len(metadata) if metadata is not None else 'Metadata is None'}")
else:
    st.error("Azure Blob Storage connection failed. Cannot load vector DB. File learning/search will be limited.")
    print("CRITICAL: Cannot load vector DB due to Blob client initialization failure (main section).")

@st.cache_data
def load_prompt_rules_cached():
    default_rules = """1. 제공된 '문서 내용'을 최우선으로 참고하여 답변합니다.
2. 질문에 대한 정보가 문서 내용에 명확히 없는 경우, "제공된 문서에서 관련 정보를 찾을 수 없습니다."라고 답변합니다. 추측성 답변은 피합니다.
3. 답변은 구체적이고 명확해야 하며, 가능하다면 관련 규정 번호나 절차 단계를 언급합니다.
4. 답변은 항상 한국어로 정중하게 제공합니다.
5. 계산이 필요한 경우, 정확한 계산 과정을 포함하여 답변합니다.
6. 사용자가 업로드한 파일(이미지 포함)의 내용과 질문을 연관지어 답변해야 할 경우, 해당 파일의 내용을 분석하여 답변에 활용합니다. 파일명도 함께 언급할 수 있습니다.
7. 답변은 항상 사용자의 질문 의도에 부합하도록 노력합니다.
8. 문서 내용에 여러 관련 정보가 있을 경우, 가장 중요하거나 질문과 직접적으로 관련된 정보를 중심으로 요약하여 답변합니다.
9. 안전, 품질, 규정 준수와 관련된 질문에는 특히 신중하고 정확한 정보를 제공합니다.
10. 답변은 문단으로 구분하여 가독성을 높입니다. 복잡한 내용은 필요시 목록 형태로 제시할 수 있습니다."""
    if os.path.exists(RULES_PATH_REPO):
        try:
            with open(RULES_PATH_REPO, "r", encoding="utf-8") as f: rules_content = f.read()
            print(f"Prompt rules loaded successfully from '{RULES_PATH_REPO}'.")
            return rules_content
        except Exception as e:
            # st.warning(f"Error loading '{RULES_PATH_REPO}': {e}. Using default rules defined above.") # UI 경고 최소화
            print(f"WARNING: Error loading prompt rules from '{RULES_PATH_REPO}': {e}. Using default rules defined in code.")
            return default_rules
    else:
        print(f"WARNING: Prompt rules file not found at '{RULES_PATH_REPO}'. Using default rules defined in code.")
        return default_rules
PROMPT_RULES_CONTENT = load_prompt_rules_cached()

def extract_text_from_file(uploaded_file_obj):
    ext = os.path.splitext(uploaded_file_obj.name)[1].lower()
    text_content = ""
    if ext in [".png", ".jpg", ".jpeg"]: # 이미지는 여기서 텍스트 추출 안 함
        print(f"DEBUG extract_text_from_file: Skipped image file '{uploaded_file_obj.name}' for text extraction.")
        return "" 
    try:
        uploaded_file_obj.seek(0)
        file_bytes = uploaded_file_obj.read()
        if ext == ".pdf":
            with fitz.open(stream=file_bytes, filetype="pdf") as doc: text_content = "\n".join(page.get_text() for page in doc)
        elif ext == ".docx": # 테이블 추출 개선 버전
            with io.BytesIO(file_bytes) as doc_io:
                doc = docx.Document(doc_io); full_text = []
                for para in doc.paragraphs: full_text.append(para.text)
                for table_idx, table in enumerate(doc.tables):
                    table_data_text = [f"--- Table {table_idx+1} Start ---"] # 테이블 구분자 추가
                    for row_idx, row in enumerate(table.rows):
                        row_cells_text = [cell.text.strip() for cell_idx, cell in enumerate(row.cells)]
                        table_data_text.append(" | ".join(row_cells_text)) # 셀 구분
                    table_data_text.append(f"--- Table {table_idx+1} End ---")
                    full_text.append("\n".join(table_data_text)) # 각 테이블 내용을 하나의 문자열로
                text_content = "\n\n".join(full_text) # 단락과 테이블 내용을 합침
        elif ext in (".xlsx", ".xlsm"):
            with io.BytesIO(file_bytes) as excel_io: df_dict = pd.read_excel(excel_io, sheet_name=None)
            text_content = "\n\n".join(f"--- Sheet: {name} ---\n{df.to_string(index=False)}" for name, df in df_dict.items())
        elif ext == ".csv":
            with io.BytesIO(file_bytes) as csv_io: # BytesIO 사용
                try: df = pd.read_csv(csv_io)
                except UnicodeDecodeError: # UTF-8 실패 시 CP949 시도
                    # BytesIO는 내부 포인터를 가지므로, 다시 읽으려면 새로운 BytesIO 객체를 만들거나 seek(0) 필요
                    df = pd.read_csv(io.BytesIO(file_bytes), encoding='cp949') 
                text_content = df.to_string(index=False)
        elif ext == ".pptx":
            with io.BytesIO(file_bytes) as ppt_io: prs = Presentation(ppt_io); text_content = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif ext == ".txt":
            try: text_content = file_bytes.decode('utf-8')
            except UnicodeDecodeError: 
                try: text_content = file_bytes.decode('cp949')
                except Exception: text_content = file_bytes.decode('latin-1', errors='replace') # 최후의 수단
            except Exception as e_txt: print(f"Txt decode error: {e_txt}"); text_content = ""
        else: st.warning(f"Unsupported text file type: {ext} (File: {uploaded_file_obj.name})"); return ""
    except Exception as e: st.error(f"Error extracting text from '{uploaded_file_obj.name}': {e}"); print(f"ERROR extracting text: {e}\n{traceback.format_exc()}"); return ""
    return text_content.strip()

def save_original_file_to_blob(uploaded_file_obj, _container_client, base_path="original_files"):
    if not _container_client or not uploaded_file_obj: return None
    try:
        safe_file_name = re.sub(r'[\\/*?:"<>|]', "_", uploaded_file_obj.name)
        blob_name = f"{base_path}/{datetime.now().strftime('%Y%m%d%H%M%S')}_{safe_file_name}"
        blob_client_instance = _container_client.get_blob_client(blob_name)
        uploaded_file_obj.seek(0)
        file_bytes_for_original = uploaded_file_obj.read() # 여기서 파일 내용을 다시 읽음
        with io.BytesIO(file_bytes_for_original) as data_stream:
            blob_client_instance.upload_blob(data_stream, overwrite=True, timeout=120)
        print(f"Successfully saved original file '{safe_file_name}' to Blob as '{blob_name}'"); return blob_name
    except Exception as e: print(f"ERROR saving original file to Blob ('{uploaded_file_obj.name}'): {e}"); return None

def log_openai_api_usage_to_blob(user_id, model_name, usage_object, _container_client, request_type="general_api_call"):
    if not _container_client: print(f"ERROR: Blob client None, cannot log API usage."); return False
    log_entry = {
        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"), "user_id": user_id, 
        "model_name": model_name, "request_type": request_type, 
        "prompt_tokens": getattr(usage_object, 'prompt_tokens', 0), 
        "completion_tokens": getattr(usage_object, 'completion_tokens', 0), 
        "total_tokens": getattr(usage_object, 'total_tokens', 0)
    }
    try:
        current_logs = load_data_from_blob(USAGE_LOG_BLOB_NAME, _container_client, "API usage log", default_value=[])
        if not isinstance(current_logs, list): current_logs = [] # 데이터가 리스트가 아니면 초기화
        current_logs.append(log_entry)
        if save_data_to_blob(current_logs, USAGE_LOG_BLOB_NAME, _container_client, "API usage log"):
            print(f"Successfully logged API usage for user '{user_id}'."); return True
        else:
            print(f"ERROR: Failed to save API usage log to Blob after appending new entry."); return False
    except Exception as e: print(f"GENERAL ERROR logging API usage: {e}\n{traceback.format_exc()}"); return False

def chunk_text_into_pieces(text_to_chunk, chunk_size=500): # 청크 크기는 토큰이 아닌 글자 수 기반
    if not text_to_chunk or not text_to_chunk.strip(): return [];
    chunks_list, current_buffer = [], ""
    for line in text_to_chunk.split("\n"): 
        stripped_line = line.strip()
        if not stripped_line and not current_buffer.strip(): continue 
        if len(current_buffer) + len(stripped_line) + 1 < chunk_size: 
            current_buffer += stripped_line + "\n"
        else: 
            if current_buffer.strip(): chunks_list.append(current_buffer.strip())
            current_buffer = stripped_line + "\n" 
    if current_buffer.strip(): chunks_list.append(current_buffer.strip())
    return [c for c in chunks_list if c] # 내용이 있는 청크만 반환

def get_image_description(image_bytes, image_filename, client_instance):
    if not client_instance: print("ERROR: OpenAI client not ready for image description."); return None
    print(f"DEBUG: Requesting description for image '{image_filename}'")
    try:
        base64_image = base64.b64encode(image_bytes).decode('utf-8')
        ext = os.path.splitext(image_filename)[1].lower(); 
        mime_type = "image/jpeg" if ext in [".jpg", ".jpeg"] else "image/png" if ext == ".png" else "application/octet-stream" # 기본값 변경
        vision_model = st.secrets.get("AZURE_OPENAI_DEPLOYMENT", "gpt-4-vision-preview") # secrets에 없으면 기본 모델명 사용
        
        response = client_instance.chat.completions.create(
            model=vision_model, 
            messages=[{"role": "user", "content": [
                {"type": "text", "text": f"Describe this image (filename: '{image_filename}') from a work/professional perspective. This description will be used for text-based search. Mention key objects, states, possible contexts, and any elements relevant to GMP/SOP if applicable."}, 
                {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{base64_image}" }}
            ]}], 
            max_tokens=IMAGE_DESCRIPTION_MAX_TOKENS, temperature=0.2, timeout=AZURE_OPENAI_TIMEOUT
        )
        description = response.choices[0].message.content.strip()
        print(f"DEBUG: Image description for '{image_filename}' generated (len: {len(description)}).")
        # 여기서 API 사용량 로깅 추가 가능
        return description
    except Exception as e: 
        print(f"ERROR during image description for '{image_filename}': {e}\n{traceback.format_exc()}")
        # st.error(f"이미지 설명 생성 오류: {e}") # UI 오류 최소화
        return None

def get_text_embedding(text_to_embed, client=openai_client, model=EMBEDDING_MODEL):
    if not client or not model or not text_to_embed or not text_to_embed.strip(): 
        # print("Skipping embedding for empty or invalid input.") # 너무 빈번한 로그 방지
        return None
    try: 
        response = client.embeddings.create(input=[text_to_embed], model=model, timeout=AZURE_OPENAI_TIMEOUT / 2)
        return response.data[0].embedding
    except Exception as e: 
        print(f"ERROR during single text embedding for text starting with '{text_to_embed[:30]}...': {e}")
        return None

def get_batch_embeddings(texts_to_embed, client=openai_client, model=EMBEDDING_MODEL, batch_size=EMBEDDING_BATCH_SIZE):
    if not client or not model or not texts_to_embed: return []
    all_embeddings = []
    for i in range(0, len(texts_to_embed), batch_size):
        batch = texts_to_embed[i:i + batch_size]; 
        if not batch: continue # 빈 배치면 건너뛰기
        print(f"DEBUG: Requesting embeddings for batch of {len(batch)} texts...")
        try:
            response = client.embeddings.create(input=batch, model=model, timeout=AZURE_OPENAI_TIMEOUT)
            # 응답 순서 보장을 위해 index 기준으로 정렬
            batch_embeddings = [item.embedding for item in sorted(response.data, key=lambda emb_item: emb_item.index)]
            all_embeddings.extend(batch_embeddings)
            print(f"DEBUG: Embeddings received for batch {i//batch_size + 1}.")
        except Exception as e: 
            print(f"ERROR during batch embedding for batch starting with '{batch[0][:30]}...': {e}")
            all_embeddings.extend([None] * len(batch)) # 실패 시 None으로 채움
    return all_embeddings

def search_similar_chunks(query_text, k_results=3):
    if index is None or index.ntotal == 0 or not metadata: 
        print("Vector DB empty or not loaded. Search aborted.")
        return []
    query_vector = get_text_embedding(query_text)
    if query_vector is None: 
        print("Query embedding failed. Search aborted.")
        return []
    try:
        actual_k = min(k_results, index.ntotal); 
        if actual_k == 0 : return []
        distances, indices_found = index.search(np.array([query_vector]).astype("float32"), actual_k)
        results = []
        for idx_val in indices_found[0]:
            if 0 <= idx_val < len(metadata) and isinstance(metadata[idx_val], dict):
                 results.append({
                    "source": metadata[idx_val].get("file_name", "Unknown Source"), 
                    "content": metadata[idx_val].get("content", ""), 
                    "is_image_description": metadata[idx_val].get("is_image_description", False), 
                    "original_file_extension": metadata[idx_val].get("original_file_extension", "")
                })
            else:
                print(f"Warning: Invalid index {idx_val} from FAISS search. Metadata length: {len(metadata)}")
        return results
    except Exception as e: 
        # st.error(f"Similarity search error: {e}") # UI 오류 최소화
        print(f"ERROR: Similarity search failed: {e}\n{traceback.format_exc()}"); return []

def add_document_to_vector_db_and_blob(uploaded_file_obj, processed_content_unused, text_chunks, _container_client, is_image_description=False):
    global index, metadata # 전역 변수 수정 명시
    if not text_chunks: st.warning(f"No content chunks to process for '{uploaded_file_obj.name}'."); return False
    if not _container_client: st.error("Cannot save learning results: Azure Blob client not ready."); return False
    
    file_type_log_desc = "image description" if is_image_description else "text document"
    print(f"Adding '{file_type_log_desc}' from '{uploaded_file_obj.name}' to vector DB.")
    
    chunk_embeddings = get_batch_embeddings(text_chunks)
    vectors_to_add, new_metadata_entries = [], []
    successful_embedding_count = 0

    for i, chunk in enumerate(text_chunks):
        embedding = chunk_embeddings[i] if i < len(chunk_embeddings) else None
        if embedding:
            vectors_to_add.append(embedding)
            new_metadata_entries.append({
                "file_name": uploaded_file_obj.name, "content": chunk, 
                "is_image_description": is_image_description, 
                "original_file_extension": os.path.splitext(uploaded_file_obj.name)[1].lower()
            })
            successful_embedding_count +=1
        else:
            print(f"Warning: Failed to get embedding for chunk {i+1} of '{uploaded_file_obj.name}'. Skipping.")

    if successful_embedding_count == 0: 
        st.error(f"No valid embeddings generated for '{uploaded_file_obj.name}'. Document not learned."); return False
    if successful_embedding_count < len(text_chunks):
        st.warning(f"Some content from '{uploaded_file_obj.name}' failed embedding. Only successful parts learned.")

    try:
        current_dim = np.array(vectors_to_add[0]).shape[0]
        if index is None or index.d != current_dim: 
            print(f"Re-initializing FAISS index. Old dim: {index.d if index else 'None'}, New dim: {current_dim}")
            index = faiss.IndexFlatL2(current_dim); metadata = []
        
        if vectors_to_add: index.add(np.array(vectors_to_add).astype("float32"))
        metadata.extend(new_metadata_entries)
        print(f"Added {len(vectors_to_add)} new chunks from '{uploaded_file_obj.name}'. Index total: {index.ntotal}, Dim: {index.d}")

        # FAISS 인덱스 및 메타데이터 Blob에 저장
        with tempfile.TemporaryDirectory() as tmpdir:
            temp_index_path = os.path.join(tmpdir, "temp.index")
            if index.ntotal > 0: 
                 faiss.write_index(index, temp_index_path)
                 if not save_binary_data_to_blob(temp_index_path, INDEX_BLOB_NAME, _container_client, "vector index"):
                     st.error("Failed to save vector index to Blob."); return False # 심각한 오류로 간주
            else: print(f"Skipping saving empty index to Blob: {INDEX_BLOB_NAME}")
        
        if not save_data_to_blob(metadata, METADATA_BLOB_NAME, _container_client, "metadata"):
            st.error("Failed to save metadata to Blob."); return False # 심각한 오류로 간주

        # 업로드 로그 기록
        uploader_name = st.session_state.user.get("name", "N/A")
        log_entry = {"file": uploaded_file_obj.name, "type": file_type_log_desc, "time": datetime.now().strftime("%Y-%m-%d %H:%M:%S"), "chunks_added": len(vectors_to_add), "uploader": uploader_name}
        upload_logs = load_data_from_blob(UPLOAD_LOG_BLOB_NAME, _container_client, "upload log", default_value=[])
        if not isinstance(upload_logs, list): upload_logs = []
        upload_logs.append(log_entry)
        if not save_data_to_blob(upload_logs, UPLOAD_LOG_BLOB_NAME, _container_client, "upload log"):
            st.warning("Failed to save upload log to Blob.") # 경고만 표시
        return True
    except Exception as e: 
        st.error(f"Error during document learning or Azure Blob upload for '{uploaded_file_obj.name}': {e}")
        print(f"ERROR: Failed to add document or upload to Blob: {e}\n{traceback.format_exc()}"); return False


# --- 탭 정의 ---
chat_interface_tab, admin_settings_tab = None, None
# current_user_info는 로그인 성공 후 정의되므로, 탭 정의는 그 이후 또는 여기서 조건부로 가능
if st.session_state.authenticated and current_user_info.get("role") == "admin":
    chat_interface_tab, admin_settings_tab = st.tabs(["💬 챗봇 질문", "⚙️ 관리자 설정"])
else: # 일반 사용자 또는 아직 current_user_info가 없을 경우 (로그인 화면)
    chat_interface_tab = st.container() 

# --- 챗봇 질문 인터페이스 ---
if chat_interface_tab: # 이 탭이 활성화되었거나, 일반 사용자의 경우 항상 이 블록 실행
    with chat_interface_tab:
        st.header("업무 질문")
        st.markdown("💡 예시: SOP 백업 주기, PIC/S Annex 11 차이, (파일 첨부 후) 이 사진 속 상황은 어떤 규정에 해당하나요? 등")

        # 채팅 메시지 표시 (current_chat_messages 사용)
        for msg_item in st.session_state.current_chat_messages:
            role, content, time_str = msg_item.get("role"), msg_item.get("content", ""), msg_item.get("time", "")
            align_class = "user-align" if role == "user" else "assistant-align"
            bubble_class = "user-bubble" if role == "user" else "assistant-bubble"
            st.markdown(f"""<div class="chat-bubble-container {align_class}"><div class="bubble {bubble_class}">{content}</div><div class="timestamp">{time_str}</div></div>""", unsafe_allow_html=True)

        st.markdown("<div style='height:16px'></div>", unsafe_allow_html=True) 
        
        # 파일 업로더 토글 버튼
        if st.button("📂 파일 첨부/숨기기", key="toggle_chat_uploader_v7_del"): 
            st.session_state.show_uploader = not st.session_state.get("show_uploader", False)

        chat_file_uploader_key = "chat_file_uploader_v7_del_widget" 
        uploaded_chat_file_runtime = None # 현재 실행에서 사용자가 업로드한 파일
        
        if st.session_state.get("show_uploader", False): # 업로더 표시 상태이면
            uploaded_chat_file_runtime = st.file_uploader("질문과 함께 참고할 파일 첨부 (선택 사항)",
                                     type=["pdf","docx","xlsx","xlsm","csv","pptx", "txt", "png", "jpg", "jpeg"], 
                                     key=chat_file_uploader_key)
            if uploaded_chat_file_runtime: 
                st.caption(f"첨부됨: {uploaded_chat_file_runtime.name} ({uploaded_chat_file_runtime.type}, {uploaded_chat_file_runtime.size} bytes)")
                if uploaded_chat_file_runtime.type.startswith("image/"): st.image(uploaded_chat_file_runtime, width=200)

        # 채팅 입력 폼
        with st.form("chat_input_form_v7_del", clear_on_submit=True): 
            query_input_col, send_button_col = st.columns([4,1])
            with query_input_col:
                user_query_input_form = st.text_input("질문 입력:", placeholder="여기에 질문을 입력하세요...", key="user_query_text_input_v7_del", label_visibility="collapsed") 
            with send_button_col: send_query_button_form = st.form_submit_button("전송")

        if send_query_button_form and user_query_input_form.strip(): # 전송 버튼 눌리고 내용 있으면
            if not openai_client or not tokenizer: # 필수 클라이언트/라이브러리 확인
                st.error("OpenAI 서비스 또는 토크나이저가 준비되지 않아 답변을 생성할 수 없습니다. 관리자에게 문의하세요."); st.stop()
            
            timestamp_now_str = datetime.now().strftime("%Y-%m-%d %H:%M")
            
            # 사용자 메시지 구성 (첨부 파일 정보 포함)
            user_message_content_for_display = user_query_input_form
            if uploaded_chat_file_runtime: 
                user_message_content_for_display += f"\n(첨부 파일: {uploaded_chat_file_runtime.name})"
            
            # 현재 채팅 메시지 목록에 사용자 메시지 추가
            st.session_state.current_chat_messages.append({"role":"user", "content":user_message_content_for_display, "time":timestamp_now_str})
            
            # 만약 현재 새 대화 상태였다면(active_conversation_id is None), 
            # 이 시점에서 대화 ID를 생성하고 all_user_conversations에 임시로 추가해둘 수 있음.
            # 또는 archive_current_chat_session_if_needed()가 호출될 때 처리되도록 함.
            # 여기서는 사용자가 메시지를 보낸 후, AI 답변까지 받고 나서 컨텍스트 전환 시 저장하는 것으로 간주.
            
            user_name_for_log = current_user_info.get("name", "anonymous_chat_user")
            print(f"User '{user_name_for_log}' submitted query: '{user_query_input_form[:50]}...' (File: {uploaded_chat_file_runtime.name if uploaded_chat_file_runtime else 'None'})")
            
            with st.spinner("답변 생성 중... 잠시만 기다려주세요."):
                assistant_response_content = "답변 생성 중 오류가 발생했습니다. 다시 시도해주세요." # 기본 오류 메시지
                try: 
                    print("Step 1: Preparing context and calculating tokens...")
                    context_items_for_llm_prompt = [] # LLM 프롬프트에 포함될 컨텍스트 아이템
                    
                    # 채팅 시 첨부 파일 처리
                    text_content_from_chat_file, is_chat_file_image, chat_file_source_display_name = None, False, None
                    if uploaded_chat_file_runtime:
                        file_extension_chat = os.path.splitext(uploaded_chat_file_runtime.name)[1].lower()
                        is_chat_file_image = file_extension_chat in [".png", ".jpg", ".jpeg"]
                        
                        if is_chat_file_image:
                            print(f"DEBUG Chat: Processing uploaded image '{uploaded_chat_file_runtime.name}' for description.")
                            with st.spinner(f"첨부 이미지 '{uploaded_chat_file_runtime.name}' 분석 중..."):
                                image_bytes = uploaded_chat_file_runtime.getvalue()
                                description = get_image_description(image_bytes, uploaded_chat_file_runtime.name, openai_client)
                            if description:
                                text_content_from_chat_file = description
                                chat_file_source_display_name = f"사용자 첨부 이미지: {uploaded_chat_file_runtime.name}"
                                print(f"DEBUG Chat: Image description generated (len: {len(description)}).")
                            else: st.warning(f"이미지 '{uploaded_chat_file_runtime.name}' 설명을 생성하지 못했습니다.")
                        else: 
                            print(f"DEBUG Chat: Extracting text from uploaded file '{uploaded_chat_file_runtime.name}'.")
                            text_content_from_chat_file = extract_text_from_file(uploaded_chat_file_runtime)
                            if text_content_from_chat_file: 
                                chat_file_source_display_name = f"사용자 첨부 파일: {uploaded_chat_file_runtime.name}"
                                print(f"DEBUG Chat: Text extracted (len: {len(text_content_from_chat_file)}).")
                            elif text_content_from_chat_file == "": st.info(f"파일 '{uploaded_chat_file_runtime.name}'이 비었거나 내용을 추출할 수 없습니다.")
                        
                        if text_content_from_chat_file: 
                            context_items_for_llm_prompt.append({
                                "source": chat_file_source_display_name, "content": text_content_from_chat_file, 
                                "is_image_description": is_chat_file_image
                            })
                    
                    # 프롬프트 구성 및 토큰 계산
                    prompt_template_for_llm = f"{PROMPT_RULES_CONTENT}\n\n다음은 사용자의 질문에 답변하는 데 도움이 되는 문서 내용입니다:\n<문서 시작>\n{{context}}\n<문서 끝>"
                    base_prompt_tokens = len(tokenizer.encode(prompt_template_for_llm.replace('{context}', '')))
                    user_query_tokens = len(tokenizer.encode(user_query_input_form))
                    max_context_tokens_allowed = TARGET_INPUT_TOKENS_FOR_PROMPT - base_prompt_tokens - user_query_tokens
                    
                    final_context_string_for_llm = "현재 참고할 수 있는 문서가 없습니다."
                    if max_context_tokens_allowed > 0:
                        query_for_vector_db_search = user_query_input_form
                        if is_chat_file_image and text_content_from_chat_file: # 이미지 설명이 있으면 검색 쿼리에 추가
                            query_for_vector_db_search = f"{user_query_input_form}\n\n첨부 이미지 내용: {text_content_from_chat_file}"
                        
                        retrieved_db_chunks = search_similar_chunks(query_for_vector_db_search, k_results=3)
                        if retrieved_db_chunks: context_items_for_llm_prompt.extend(retrieved_db_chunks)
                        
                        if context_items_for_llm_prompt:
                            unique_contents_seen = set()
                            formatted_context_segments = []
                            for item in context_items_for_llm_prompt:
                                content_segment = item.get("content","").strip()
                                if content_segment and content_segment not in unique_contents_seen:
                                    source_name = item.get('source','알 수 없음').replace("사용자 첨부 이미지: ","").replace("사용자 첨부 파일: ","")
                                    prefix = "[이미지 설명: " if item.get("is_image_description") else "[출처 문서: "
                                    formatted_context_segments.append(f"{prefix}{source_name}]\n{content_segment}")
                                    unique_contents_seen.add(content_segment)
                            
                            if formatted_context_segments:
                                combined_context_str = "\n\n---\n\n".join(formatted_context_segments)
                                encoded_combined_context = tokenizer.encode(combined_context_str)
                                if len(encoded_combined_context) > max_context_tokens_allowed:
                                    truncated_tokens_for_context = encoded_combined_context[:max_context_tokens_allowed]
                                    final_context_string_for_llm = tokenizer.decode(truncated_tokens_for_context)
                                    if len(encoded_combined_context) > len(truncated_tokens_for_context):
                                        final_context_string_for_llm += "\n(...문서 내용이 길어 일부 잘렸을 수 있습니다.)"
                                else: final_context_string_for_llm = combined_context_str
                    
                    system_prompt_final = prompt_template_for_llm.replace('{context}', final_context_string_for_llm)
                    total_input_tokens = len(tokenizer.encode(system_prompt_final)) + user_query_tokens
                    if total_input_tokens > MODEL_MAX_INPUT_TOKENS: 
                        print(f"CRITICAL WARNING: Total input tokens ({total_input_tokens}) exceed model max ({MODEL_MAX_INPUT_TOKENS})!")
                    
                    api_messages_to_send = [{"role":"system", "content": system_prompt_final}, {"role":"user", "content": user_query_input_form}]
                    print("Step 2: Sending request to Azure OpenAI for chat completion...")
                    
                    chat_model_deployment_name = st.secrets.get("AZURE_OPENAI_DEPLOYMENT")
                    if not chat_model_deployment_name:
                        st.error("채팅 모델 배포 이름('AZURE_OPENAI_DEPLOYMENT')이 secrets에 없습니다."); raise ValueError("Chat model name missing.")
                        
                    chat_completion_result = openai_client.chat.completions.create(
                        model=chat_model_deployment_name, messages=api_messages_to_send,
                        max_tokens=MODEL_MAX_OUTPUT_TOKENS, temperature=0.1, timeout=AZURE_OPENAI_TIMEOUT
                    )
                    assistant_response_content = chat_completion_result.choices[0].message.content.strip()
                    print("Azure OpenAI response received.")

                    if chat_completion_result.usage and container_client:
                        log_openai_api_usage_to_blob(user_name_for_log, chat_model_deployment_name, chat_completion_result.usage, container_client, request_type="chat_completion_with_rag")
                
                except Exception as gen_err: 
                    assistant_response_content = f"답변 생성 중 예상치 못한 오류 발생: {gen_err}."
                    st.error(assistant_response_content) # UI에 오류 표시
                    print(f"UNEXPECTED ERROR during response generation: {gen_err}\n{traceback.format_exc()}")

            st.session_state.current_chat_messages.append({"role":"assistant", "content":assistant_response_content, "time":timestamp_now_str})
            # 답변 후 현재 대화가 새 대화였으면 ID를 부여하고 저장할 준비 (실제 저장은 컨텍스트 전환 시)
            if st.session_state.active_conversation_id is None and st.session_state.current_chat_messages:
                # archive_current_chat_session_if_needed() 함수가 호출될 때 새 ID가 할당되고 저장됨.
                # 사용자가 메시지를 계속 보내는 동안은 active_id는 None으로 유지될 수 있으며,
                # 새 대화 시작 / 다른 대화 로드 / 로그아웃 시 아카이브 함수가 호출되면서 ID가 부여되고 저장됨.
                # 만약, 첫 응답 직후 바로 all_user_conversations에 반영하고 싶다면 여기서 archive 함수 호출 필요.
                # 지금은 archive 함수가 컨텍스트 전환 시점에 호출되도록 되어 있으므로, 여기서는 특별한 처리 안함.
                pass
            
            print("Response processing complete. Triggering rerun to display new messages."); st.rerun()

# --- 관리자 설정 탭 ---
if admin_settings_tab: # admin_settings_tab이 None이 아니고, 현재 활성화된 탭일 때 (st.tabs 사용 시 자동 처리)
    with admin_settings_tab:
        st.header("⚙️ 관리자 설정")
        # 가입 승인 대기자
        st.subheader("👥 가입 승인 대기자")
        if not USERS or not isinstance(USERS, dict): st.warning("사용자 정보를 로드할 수 없거나 형식이 올바르지 않습니다.")
        else:
            pending_users = {uid:udata for uid,udata in USERS.items() if isinstance(udata, dict) and not udata.get("approved")}
            if pending_users:
                for pending_uid, pending_data in pending_users.items():
                    with st.expander(f"{pending_data.get('name','N/A')} ({pending_uid}) - {pending_data.get('department','N/A')}"):
                        approve_col, reject_col = st.columns(2)
                        if approve_col.button("승인", key=f"admin_approve_user_v7_{pending_uid}"): 
                            USERS[pending_uid]["approved"] = True
                            if save_data_to_blob(USERS, USERS_BLOB_NAME, container_client, "user info (approval)"):
                                st.success(f"사용자 '{pending_uid}' 승인 완료."); st.rerun()
                            else: st.error("사용자 승인 정보 저장 실패.")
                        if reject_col.button("거절", key=f"admin_reject_user_v7_{pending_uid}"): 
                            USERS.pop(pending_uid, None)
                            if save_data_to_blob(USERS, USERS_BLOB_NAME, container_client, "user info (rejection)"):
                                st.info(f"사용자 '{pending_uid}' 거절 처리 완료."); st.rerun()
                            else: st.error("사용자 거절 정보 저장 실패.")
            else: st.info("승인 대기 중인 사용자가 없습니다.")
        st.markdown("---")

        # 파일 업로드 및 학습
        st.subheader("📁 파일 업로드 및 학습 (Azure Blob Storage)")
        if 'processed_admin_file_info' not in st.session_state: st.session_state.processed_admin_file_info = None
        def clear_processed_admin_file_info_callback(): st.session_state.processed_admin_file_info = None
        
        admin_uploaded_file_widget = st.file_uploader(
            "학습할 파일 업로드 (PDF, DOCX, XLSX, CSV, PPTX, TXT, PNG, JPG, JPEG)",
            type=["pdf","docx","xlsx","xlsm","csv","pptx", "txt", "png", "jpg", "jpeg"], 
            key="admin_file_uploader_v7_del",
            on_change=clear_processed_admin_file_info_callback,
            accept_multiple_files=False 
        )

        if admin_uploaded_file_widget and container_client:
            current_admin_file_details = (admin_uploaded_file_widget.name, admin_uploaded_file_widget.size, admin_uploaded_file_widget.type)
            if st.session_state.processed_admin_file_info != current_admin_file_details: # 중복 처리 방지
                print(f"DEBUG Admin Upload: New file detected by admin. Info: {current_admin_file_details}")
                try: 
                    file_ext_admin_ul = os.path.splitext(admin_uploaded_file_widget.name)[1].lower()
                    is_admin_upload_image = file_ext_admin_ul in [".png", ".jpg", ".jpeg"]
                    content_to_learn, is_description_for_learning = None, False

                    if is_admin_upload_image:
                        with st.spinner(f"이미지 '{admin_uploaded_file_widget.name}' 처리 및 설명 생성 중..."):
                            admin_img_bytes = admin_uploaded_file_widget.getvalue()
                            admin_img_description = get_image_description(admin_img_bytes, admin_uploaded_file_widget.name, openai_client)
                        if admin_img_description:
                            content_to_learn = admin_img_description; is_description_for_learning = True
                            st.info(f"이미지 '{admin_uploaded_file_widget.name}' 설명 생성 (길이: {len(admin_img_description)}). 이 설명이 학습됩니다.")
                            st.text_area("생성된 이미지 설명 (학습용)", admin_img_description, height=150, disabled=True)
                        else: st.error(f"이미지 '{admin_uploaded_file_widget.name}' 설명 생성 실패. 학습 제외.")
                    else: 
                        with st.spinner(f"'{admin_uploaded_file_widget.name}'에서 텍스트 추출 중..."):
                            content_to_learn = extract_text_from_file(admin_uploaded_file_widget)
                        if content_to_learn: st.info(f"'{admin_uploaded_file_widget.name}' 텍스트 추출 (길이: {len(content_to_learn)}).")
                        else: st.warning(f"'{admin_uploaded_file_widget.name}' 내용 추출 불가 또는 비어있음. 학습 제외.")
                    
                    if content_to_learn: 
                        with st.spinner(f"'{admin_uploaded_file_widget.name}' 내용 처리 및 학습 중..."):
                            chunks_for_learning = chunk_text_into_pieces(content_to_learn)
                            if chunks_for_learning:
                                original_blob_path = save_original_file_to_blob(admin_uploaded_file_widget, container_client)
                                if original_blob_path: st.caption(f"원본 파일 '{admin_uploaded_file_widget.name}' Blob 저장: '{original_blob_path}'.")
                                else: st.warning(f"원본 파일 '{admin_uploaded_file_widget.name}' Blob 저장 실패.")

                                if add_document_to_vector_db_and_blob(admin_uploaded_file_widget, content_to_learn, chunks_for_learning, container_client, is_image_description=is_description_for_learning):
                                    st.success(f"파일 '{admin_uploaded_file_widget.name}' 학습 및 Azure Blob Storage 업데이트 완료!")
                                    st.session_state.processed_admin_file_info = current_admin_file_details 
                                    st.rerun() 
                                else: st.error(f"'{admin_uploaded_file_widget.name}' 학습 또는 Blob 업데이트 중 오류."); st.session_state.processed_admin_file_info = None 
                            else: st.warning(f"'{admin_uploaded_file_widget.name}'에 대한 학습 청크 생성 안됨."); st.session_state.processed_admin_file_info = None
                except Exception as e_admin_file_main_proc:
                    st.error(f"관리자 업로드 파일 {admin_uploaded_file_widget.name} 처리 중 오류: {e_admin_file_main_proc}")
                    print(f"CRITICAL ERROR in admin_upload_processing for {admin_uploaded_file_widget.name}: {e_admin_file_main_proc}\n{traceback.format_exc()}")
                    st.session_state.processed_admin_file_info = None
            elif st.session_state.processed_admin_file_info == current_admin_file_details:
                 st.caption(f"파일 '{admin_uploaded_file_widget.name}'은 이전에 처리되었습니다. 재학습하려면 다시 업로드하세요.")
        elif admin_uploaded_file_widget and not container_client:
            st.error("파일 업로드 및 학습 불가: Azure Blob 클라이언트가 준비되지 않았습니다.")
        st.markdown("---")

        # API 사용량 모니터링
        st.subheader("📊 API 사용량 모니터링 (Blob 로그 기반)")
        if container_client:
            usage_data_list = load_data_from_blob(USAGE_LOG_BLOB_NAME, container_client, "API usage log", default_value=[])
            if usage_data_list and isinstance(usage_data_list, list) and len(usage_data_list) > 0 :
                df_usage = pd.DataFrame(usage_data_list)
                for col_token in ["total_tokens", "prompt_tokens", "completion_tokens"]: 
                    df_usage[col_token] = pd.to_numeric(df_usage.get(col_token, 0), errors='coerce').fillna(0)
                if "request_type" not in df_usage.columns: df_usage["request_type"] = "unknown"
                
                total_tokens_all = df_usage["total_tokens"].sum()
                st.metric("총 API 호출 수", len(df_usage))
                st.metric("총 사용 토큰 수", f"{int(total_tokens_all):,}")

                token_cost_config = 0.0
                try: token_cost_config = float(st.secrets.get("TOKEN_COST","0.0"))
                except (ValueError, TypeError): pass 
                st.metric("예상 비용 (USD)", f"${total_tokens_all * token_cost_config:.4f}") 

                if "timestamp" in df_usage.columns:
                    try: 
                         df_usage['timestamp'] = pd.to_datetime(df_usage['timestamp'])
                         st.dataframe(df_usage.sort_values(by="timestamp",ascending=False), use_container_width=True, hide_index=True)
                    except Exception as e_sort_usage_df:
                         print(f"Warning: Could not sort usage log by timestamp: {e_sort_usage_df}")
                         st.dataframe(df_usage, use_container_width=True, hide_index=True) 
                else: st.dataframe(df_usage, use_container_width=True, hide_index=True)
            else: st.info("Blob에 기록된 API 사용량 데이터가 없거나 비어있습니다.")
        else: st.warning("API 사용량 모니터링 표시 불가: Azure Blob 클라이언트가 준비되지 않았습니다.")
        st.markdown("---")

        # Azure Blob Storage 파일 목록
        st.subheader("📂 Azure Blob Storage 파일 목록 (최근 100개)")
        if container_client:
            try:
                blob_display_list = []
                blob_items_sorted = sorted(container_client.list_blobs(), key=lambda b: b.last_modified, reverse=True)
                for count_blob, blob_item_data in enumerate(blob_items_sorted):
                    if count_blob >= 100: break
                    blob_display_list.append({
                        "파일명": blob_item_data.name, "크기 (bytes)": blob_item_data.size,
                        "수정일": blob_item_data.last_modified.strftime('%Y-%m-%d %H:%M:%S') if blob_item_data.last_modified else 'N/A'
                    })
                if blob_display_list: st.dataframe(pd.DataFrame(blob_display_list), use_container_width=True, hide_index=True)
                else: st.info("Azure Blob Storage에 파일이 없습니다.")
            except AzureError as ae_blob_list: 
                 st.error(f"Azure Blob 파일 목록 조회 중 Azure 서비스 오류: {ae_blob_list}")
                 print(f"AZURE ERROR listing blobs: {ae_blob_list}\n{traceback.format_exc()}")
            except Exception as e_blob_list:
                st.error(f"Azure Blob 파일 목록 조회 중 알 수 없는 오류: {e_blob_list}")
                print(f"ERROR listing blobs: {e_blob_list}\n{traceback.format_exc()}")
        else: st.warning("파일 목록 표시 불가: Azure Blob 클라이언트가 준비되지 않았습니다.")
